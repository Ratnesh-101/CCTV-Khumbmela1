"""
Generate hackathon deck (technical + Kumbh Mela use case + code pipeline).
Run: python scripts/generate_hackathon_pptx.py
Requires: pip install python-pptx
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
# Default under presentations/; root copy optional (often locked if PowerPoint has it open).
OUT = ROOT / "presentations" / "CCTV_Risk_Hackathon_Deck.pptx"
OUT_ROOT = ROOT / "CCTV_Risk_Hackathon_Deck.pptx"
FONT = Pt(17)
FONT_SM = Pt(15)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str], small: bool = False) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    fs = FONT_SM if small else FONT
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = fs


def main() -> None:
    prs = Presentation()

    _add_title_slide(
        prs,
        "CCTV Risk Intelligence Stack",
        "Hackathon prototype — Kumbh Mela–scale fixed CCTV: fuse perception into one risk trend + operator aids\n"
        "Python · YOLOv8 · OpenCV · PyTorch · Streamlit · LangChain · FAISS · scikit-learn · Gymnasium",
    )

    _add_bullet_slide(
        prs,
        "Kumbh Mela — why this example fits",
        [
            "Massive pedestrian density on ghats, bridges, corridors; fixed cameras are standard ops tech.",
            "Control rooms cannot watch every tile at equal attention — need prioritization signals.",
            "We target scene-level risk (crowding, motion, unattended objects, incident cues) — not pilgrim ID.",
            "Complements human judgment: ranked feeds + SOP text (RAG), not autonomous enforcement.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Issue the pipeline solves",
        [
            "Fragmentation: boxes, motion, and rules are separate — operators lack one fused trend.",
            "Solution: transparent 0–100 score + named flags (e.g. disturbance_motion, dense_crowd) per frame.",
            "Escalation structure: tabular Q-learning suggests discrete action with cooldown debouncing.",
            "Procedure recall: RAG retrieves playbook chunks (crowd, lost property, evacuation language).",
            "Kumbh link: proximity + motion ≈ corridor pressure; bags ≈ security SOP; optional audio ≈ incident cue (demo).",
        ],
    )

    _add_bullet_slide(
        prs,
        "Code pipeline — repository map",
        [
            "app.py — Streamlit: upload video, loop, Plotly charts, st.video.",
            "process_video.py — batch: annotated.mp4 + events.csv.",
            "src/video_pipeline.py — VideoAnalyzer.analyze_frame → FrameSignals; draw_overlay().",
            "src/risk_scorer.py — compute_risk(), WEIGHTS, semantic labels.",
            "src/rl_policy.py — rl_action(), models/q_table.npy.",
            "src/rag_ops.py — retrieve_context(), rag_operator_brief().",
            "src/video_io.py — open_video_writer() H.264 / MJPG fallback.",
            "src/wifi_drone_anomaly.py — Isolation Forest + optional embeddings.",
        ],
        small=True,
    )

    _add_bullet_slide(
        prs,
        "Core loop (same pattern as app.py)",
        [
            "sig = analyzer.analyze_frame(frame, prev_gray, gray, tracker)",
            "rs = compute_risk(person_count, motion_score, proximity, abandoned_sec,",
            "    weapon_like, fight_like, theft_like, accident_like, missile, drone, wifi, …)",
            "act_id, act_name = rl_action(rs.total, cooldown, Q_table)",
            "writer.write(draw_overlay(frame, sig, rs.total, rs.labels))",
        ],
        small=True,
    )

    _add_bullet_slide(
        prs,
        "Perception & fusion",
        [
            "YOLOv8: persons, bags, custom classes (fight, accident, weapon…) via runtime name sets.",
            "Farneback optical flow on person-masked ROIs → motion_score.",
            "Centroid proximity on persons → crowding; SimpleBagTracker → abandoned seconds.",
            "Weighted linear fusion Σ w_i c_i × 100 → total; thresholded labels for UI strip.",
        ],
    )

    _add_bullet_slide(
        prs,
        "RL · RAG · optional sensors",
        [
            "risk_bucket(score) → discrete state for Q-table; debounced escalation hint.",
            "LangChain + FAISS + sentence-transformers for operator playbook queries.",
            "Wi-Fi tab: IF on flow features — optional scalar fused into video risk.",
            "Audio tab: transient heuristic → acoustic_threat (demo; not gunshot classifier).",
        ],
    )

    _add_bullet_slide(
        prs,
        "Control-room outputs",
        [
            "Annotated video: boxes, [class tags], green risk-flag band.",
            "Plotly: risk vs time + stacked normalized components.",
            "CSV / table: per-frame audit trail; optional fixed camera lat/lng on alerts.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Limits @ Mela-scale & ethics",
        [
            "Not certified crowd-safety; heuristics false-alarm — human-in-the-loop mandatory.",
            "Train YOLO on domain CCTV for fights / accidents / crush proxies; COCO is insufficient alone.",
            "Governance: public-safety ops only; comply with surveillance law; avoid mass individual tracking.",
            "Next: edge ONNX/TensorRT, multi-camera sync, sector-specific thresholds, incident workflow.",
            "Q&A — thank you.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("Wrote", OUT.resolve(), "(9 slides)")
    try:
        shutil.copy2(OUT, OUT_ROOT)
        print("Also copied to", OUT_ROOT.resolve())
    except OSError as e:
        print("Skip root copy (close PowerPoint if file is open):", e)


if __name__ == "__main__":
    main()
